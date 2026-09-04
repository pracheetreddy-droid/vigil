import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Dark Cyber Theme
    BG_COLOR = RGBColor(9, 9, 13)       # #09090d
    CARD_BG = RGBColor(18, 18, 26)       # #12121a
    TEXT_WHITE = RGBColor(244, 244, 245) # #f4f4f5
    TEXT_MUTED = RGBColor(161, 161, 170) # #a1a1aa
    BRAND_INDIGO = RGBColor(99, 102, 241)# #6366f1
    ACCENT_CYAN = RGBColor(6, 182, 212)  # #06b6d4
    EMERALD_GREEN = RGBColor(16, 185, 129)# #10b981
    ALERT_RED = RGBColor(239, 68, 68)    # #ef4444

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="NEXORA ROUND 2 • VIGIL"):
        # Category label
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BRAND_INDIGO

        # Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text.upper()
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "NEXORA • THE VIBE-CODING CLUB (NITTE)"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = BRAND_INDIGO
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "VIGIL"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "ALWAYS AWARE. ALWAYS WITH YOU."
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_CYAN
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "AI-Powered Tourist Safety & Location Intelligence Platform\nCandidate: Pracheet | Track: Tourist Safety Platform"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.alignment = PP_ALIGN.CENTER

    # SLIDE 2: Problem Statement
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Problem Statement — The Challenge in Tourist Safety")

    probs = [
        ("1. Navigational Blindness", "Tourists unfamiliar with local cities unknowingly wander into dimly lit side alleys and high-risk crime zones after hours."),
        ("2. High Emergency Response Friction", "When emergencies occur, travelers panic and struggle to transmit accurate GPS coordinates or cryptographic identity data to local police."),
        ("3. Authority Information Asymmetry", "State law enforcement and emergency response centers lack predictive spatio-temporal risk models to deploy patrols before crime surges occur.")
    ]

    for idx, (title, desc) in enumerate(probs):
        box = slide2.shapes.add_textbox(Inches(0.8 + idx * 3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ALERT_RED if idx == 0 else BRAND_INDIGO if idx == 1 else ACCENT_CYAN

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED

    # SLIDE 3: Solution Overview
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. Solution Overview — The VIGIL Safety Ecosystem")

    txBox = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "PRODUCT VISION ARCHITECTURE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BRAND_INDIGO

    p_flow = tf.add_paragraph()
    p_flow.text = "TOURIST ➔ LOCATION INTELLIGENCE ➔ AI RISK ANALYSIS ➔ EMERGENCY RESPONSE ➔ AUTHORITIES"
    p_flow.font.size = Pt(16)
    p_flow.font.bold = True
    p_flow.font.color.rgb = EMERALD_GREEN

    points = [
        ("• Tourist Safety Suite", "Equips travelers with real-time sector risk heatmaps, transparent 0-100 safety scoring, QR digital identity, and 1-touch SOS."),
        ("• VIGIL SafeRoute Navigation", "Optimizes journey paths by balancing transit duration against crime heatmaps & CCTV lighting corridors."),
        ("• VIGIL COMMAND Authority Portal", "Enables law enforcement to monitor live SOS queues, inspect verified tourist credentials, and triage incidents."),
        ("• AI Hotspot Intelligence", "Predicts emerging crime hotspots (+34% surge analysis) and recommends targeted patrol deployment.")
    ]

    for p_title, p_desc in points:
        pt = tf.add_paragraph()
        pt.text = f"\n{p_title}"
        pt.font.size = Pt(15)
        pt.font.bold = True
        pt.font.color.rgb = TEXT_WHITE

        pd = tf.add_paragraph()
        pd.text = p_desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_MUTED

    # SLIDE 4: Main Features & UX
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Core Features & User Experience")

    feats = [
        ("Interactive Safety Map", "Dark tactical Leaflet canvas showing 🟢 Low, 🟡 Moderate, 🔴 High risk sectors and clickable incident drawers."),
        ("VIGIL Safety Index (87/100)", "Transparent scoring engine factoring incident density (18%), emergency access (92%), recent activity (12%), and time risk (20%)."),
        ("QR Digital Tourist ID", "Cryptographic QR identity card (/id & /verify/VG-284921) verified by Karnataka Tourism Authority."),
        ("1-Touch SOS Workflow", "3-second countdown hold preventing accidental triggers, broadcasting live GPS telemetry to Authority Command.")
    ]

    for idx, (ft_title, ft_desc) in enumerate(feats):
        col = idx % 2
        row = idx // 2
        box = slide4.shapes.add_textbox(Inches(0.8 + col * 5.9), Inches(1.8 + row * 2.6), Inches(5.6), Inches(2.3))
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ft_title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN

        pd = tf.add_paragraph()
        pd.text = ft_desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_MUTED

    # SLIDE 5: AI & Smart Features
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. AI & Smart Features — 3 Key Innovations")

    ais = [
        ("1. VIGIL AI Safety Copilot", "Context-aware AI assistant understanding live GPS position, safety index, and hospital distances. Outputs structured risk advice."),
        ("2. SafeRoute Routing Engine", "Safety-aware navigation engine comparing Fastest (18m/48), Balanced (21m/71), and Safest (24m/94 - avoiding 2 risk zones)."),
        ("3. AI Hotspot Intelligence", "Predictive spatio-temporal cluster analysis detecting +34% incident spikes and recommending mobile patrol deployment.")
    ]

    for idx, (title, desc) in enumerate(ais):
        box = slide5.shapes.add_textbox(Inches(0.8 + idx * 3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = BRAND_INDIGO

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_MUTED

    # SLIDE 6: Authority Command Center
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. Authority Operations — VIGIL COMMAND (/authority)")

    txBox = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    items = [
        ("• Real-Time Metrics Overview", "Tracks 03 Active SOS alerts, 27 Incidents, 1,284 Monitored Tourists, and 08 High-Risk Sectors."),
        ("• Active SOS Dispatch Queue", "Allows emergency dispatchers to view live GPS coordinates, call tourists directly, or inspect verified credentials."),
        ("• Tourist Profile Inspector", "Displays passport hash, emergency contacts, journey history, and last known check-in status."),
        ("• Recharts Analytics Suite", "Visualizes 24-hour incident density trends, category distribution (Theft, Harassment), and average response time (4.2 mins).")
    ]

    for title, desc in items:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ALERT_RED if "SOS" in title else EMERALD_GREEN

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_MUTED

    # SLIDE 7: Tech Stack
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Tech Stack & Engineering Architecture")

    stacks = [
        ("Frontend & Framework", "Next.js 14 (App Router), React 18, TypeScript, Tailwind CSS"),
        ("Design System & UI", "Dark Cybersecurity Theme, Glassmorphism backdrop-blur, Framer Motion, Lucide Icons"),
        ("Data & Visualization", "Recharts (Analytics), Leaflet (Dark Geospatial Canvas), qrcode.react (QR Identity)"),
        ("AI & Web Audio API", "Gemini 1.5 Flash API + Local Contextual AI Engine Fallback, Procedural Audio Synth")
    ]

    for idx, (category, details) in enumerate(stacks):
        box = slide7.shapes.add_textbox(Inches(0.8 + (idx % 2) * 5.9), Inches(1.8 + (idx // 2) * 2.5), Inches(5.6), Inches(2.2))
        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = BRAND_INDIGO

        pd = tf.add_paragraph()
        pd.text = details
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_MUTED

    # SLIDE 8: Unique Differentiators
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Product Differentiators — Why VIGIL Stands Out")

    diffs = [
        ("• Zero-Dependency Demo Mode", "Runs 100% locally without requiring external API keys or Mapbox/Supabase infrastructure."),
        ("• Recruiter Fast-Track Mode", "Includes an on-screen Nexora requirement audit checklist and 1-click feature teleport links."),
        ("• Procedural Sound Synth", "Uses Web Audio API to synthesize real-time radar sonar pings, AI tones, and SOS alarms."),
        ("• Multilingual Safety", "Supports English, Hindi, Kannada, Tamil, and Telugu for emergency alerts and safety guidelines."),
        ("• Automated Check-In Escalation", "15-minute countdown timer that automatically alerts Authority Command if missed.")
    ]

    txBox = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    for title, desc in diffs:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(13)
        pd.font.color.rgb = TEXT_MUTED

    # SLIDE 9: Pracheet Demo Walkthrough
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "8. Guided Demo Walkthrough — The Pracheet Scenario")

    txBox = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    steps = [
        "1. Pracheet arrives in Bengaluru and generates Digital Tourist ID (/id).",
        "2. VIGIL detects location & calculates 87/100 LOW RISK Safety Score (/dashboard).",
        "3. Pracheet inspects MG Road (72 Moderate Risk) on Safety Map (/map).",
        "4. Asks VIGIL AI: 'Is it safe to visit this area tonight?' (/copilot).",
        "5. Selects SAFEST Route (24m, score 94) avoiding 2 risk zones (/saferoute).",
        "6. Triggers SOS emergency (/sos) -> VIGIL COMMAND receives alert (/authority).",
        "7. AI Hotspot Intelligence detects +34% surge -> Authority dispatches Patrol Unit #4."
    ]

    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_WHITE

    # SLIDE 10: Conclusion
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Conclusion & Submission Links")

    txBox = slide10.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "VIGIL"
    p0.font.size = Pt(40)
    p0.font.bold = True
    p0.font.color.rgb = BRAND_INDIGO
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "Always aware. Always with you."
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "\n• GitHub Repository: https://github.com/your-username/vigil\n• Local Dev URL: http://localhost:3000\n• Built for Nexora: The Vibe-Coding Club (Round 2 Task)"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_WHITE
    p2.alignment = PP_ALIGN.CENTER

    prs.save("VIGIL_Presentation.pptx")
    print("VIGIL_Presentation.pptx generated successfully!")

if __name__ == "__main__":
    create_presentation()
